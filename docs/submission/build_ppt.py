# -*- coding: utf-8 -*-
"""Build the competition deck for XH-202621 (智演 JobEvolution).

Run:  /tmp/docenv/bin/python docs/submission/build_ppt.py
Needs: python-pptx, pillow; diagrams in /tmp/je_doc (see diagrams.py) and
screenshots in /tmp/je_shots.
"""
import os
from PIL import Image
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

ROOT = "/Users/kangvcar/Documents/code/JobEvolution"
OUT = f"{ROOT}/XH-202621_参赛答辩PPT_智演JobEvolution.pptx"
DIAG = "/tmp/je_doc"
SHOT = "/tmp/je_shots"
TMP = "/tmp/je_ppt"
os.makedirs(TMP, exist_ok=True)

INK = "1d1d1f"
PAPER = "ffffff"
CREAM = "f4f1ea"
ACCENT = "c8791b"
RULE = "8c8c8c"
GREY = "e9e9e9"
DIM = "5f5f5f"
LATIN = "Helvetica Neue"
CJK = "PingFang SC"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5
PAGE = [0]


def rgb(h):
    return RGBColor.from_string(h)


def _font(run, size, color, bold, italic=False, font=None):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    run.font.name = font or LATIN
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", CJK)


def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fill)
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = rgb(line)
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    s.text_frame.text = ""
    return s


def tb(slide, x, y, w, h, paras, size=14, color=INK, bold=False, align="l", anchor="t",
       spacing=1.15, space_after=4, font=None, margin=0.05):
    """paras: str | list of (str | dict(text,size,bold,color,italic,bullet,align,after,font))."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    tf.vertical_anchor = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}[anchor]
    if isinstance(paras, str):
        paras = [paras]
    first = True
    for p in paras:
        if isinstance(p, str):
            p = {"text": p}
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[p.get("align", align)]
        para.line_spacing = spacing
        para.space_after = Pt(p.get("after", space_after))
        txt = p.get("text", "")
        if p.get("bullet"):
            txt = "▪  " + txt
        segs = p.get("runs") or [(txt, {})]
        for seg, opt in segs:
            r = para.add_run()
            r.text = seg
            _font(r, opt.get("size", p.get("size", size)), opt.get("color", p.get("color", color)),
                  opt.get("bold", p.get("bold", bold)), opt.get("italic", p.get("italic", False)),
                  opt.get("font", p.get("font", font)))
    return box


def pic(slide, path, x, y, w, h, mode="cover", border=True, anchor="top"):
    """Place an image fitted into (w,h) inches. cover = crop; contain = letterbox."""
    im = Image.open(path).convert("RGB")
    iw, ih = im.size
    target = w / h
    if mode == "cover":
        cur = iw / ih
        if cur > target:
            nw = int(ih * target)
            left = (iw - nw) // 2
            im = im.crop((left, 0, left + nw, ih))
        else:
            nh = int(iw / target)
            top = 0 if anchor == "top" else (ih - nh) // 2
            im = im.crop((0, top, iw, top + nh))
        out = f"{TMP}/{os.path.basename(path).rsplit('.', 1)[0]}_{int(w*100)}x{int(h*100)}.png"
        im.save(out)
        if border:
            rect(slide, x - 0.02, y - 0.02, w + 0.04, h + 0.04, fill=None, line=RULE, lw=0.75)
        return slide.shapes.add_picture(out, Inches(x), Inches(y), Inches(w), Inches(h))
    # contain
    cur = iw / ih
    if cur > target:
        pw, ph = w, w / cur
    else:
        pw, ph = h * cur, h
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    return slide.shapes.add_picture(path, Inches(px), Inches(py), Inches(pw), Inches(ph))


def chrome(section, title, sub=None, notes=None):
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=PAPER)
    rect(s, 0.6, 0.55, 0.14, 0.14, fill=ACCENT)
    tb(s, 0.85, 0.42, 6, 0.4, section, size=11, color=DIM)
    tb(s, 0.6, 0.72, 11.5, 0.7, title, size=26, bold=True)
    if sub:
        tb(s, 0.6, 1.32, 11.8, 0.4, sub, size=13, color=DIM)
    rect(s, 0.6, 1.72, W - 1.2, 0.012, fill=INK)
    # footer
    rect(s, 0.6, H - 0.5, W - 1.2, 0.008, fill=GREY)
    tb(s, 0.6, H - 0.45, 8, 0.3, "智演 JobEvolution · XH-202621 多源异构数据驱动岗位和能力图谱构建与动态演化分析研究", size=9, color=RULE)
    tb(s, W - 1.6, H - 0.45, 1.0, 0.3, f"{PAGE[0]:02d}", size=10, color=RULE, align="r")
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def card(slide, x, y, w, h, head, body, fill=CREAM, head_color=INK, size=12, head_size=14):
    rect(slide, x, y, w, h, fill=fill, line=None)
    rect(slide, x, y, 0.06, h, fill=ACCENT)
    tb(slide, x + 0.2, y + 0.12, w - 0.35, 0.4, head, size=head_size, bold=True, color=head_color)
    tb(slide, x + 0.2, y + 0.55, w - 0.35, h - 0.65, body, size=size, color=INK, spacing=1.2, space_after=3)


def table(slide, x, y, w, rows, col_w, size=11, head_fill=INK, row_h=0.36, first_bold=False):
    n, m = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n, m, Inches(x), Inches(y), Inches(w), Inches(row_h * n))
    t = shp.table
    tblPr = t._tbl.tblPr
    tblPr.set("bandRow", "0")
    tblPr.set("firstRow", "0")
    # remove default style id
    for el in tblPr.findall(qn("a:tableStyleId")):
        tblPr.remove(el)
    for j, cw in enumerate(col_w):
        t.columns[j].width = Inches(cw)
    for i, row in enumerate(rows):
        t.rows[i].height = Inches(row_h)
        for j, val in enumerate(row):
            c = t.cell(i, j)
            c.margin_left = c.margin_right = Inches(0.08)
            c.margin_top = c.margin_bottom = Inches(0.04)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.fill.solid()
            if i == 0:
                c.fill.fore_color.rgb = rgb(head_fill)
            else:
                c.fill.fore_color.rgb = rgb(CREAM if i % 2 == 1 else PAPER)
            tf = c.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            _font(r, size, PAPER if i == 0 else INK, bold=(i == 0) or (first_bold and j == 0))
            # borders
            tcPr = c._tc.get_or_add_tcPr()
            for k, tag in enumerate(("a:lnL", "a:lnR", "a:lnT", "a:lnB")):
                ln = etree.Element(qn(tag))
                tcPr.insert(k, ln)
                ln.set("w", "6350")
                sf = etree.SubElement(ln, qn("a:solidFill"))
                clr = etree.SubElement(sf, qn("a:srgbClr"))
                clr.set("val", RULE)
    return shp


def stat(slide, x, y, w, num, label, num_size=34, color=INK):
    tb(slide, x, y, w, 0.7, num, size=num_size, bold=True, color=color)
    tb(slide, x, y + 0.68, w, 0.5, label, size=11, color=DIM)


# =============================================================== 1 cover
def s_cover():
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=CREAM)
    rect(s, 0, 0, 0.35, H, fill=INK)
    rect(s, 0.35, 0, 0.06, H, fill=ACCENT)
    tb(s, 1.0, 0.75, 9, 0.4, "「挑战杯」揭榜挂帅专项赛 · 科大讯飞发榜 · 赛题编号 XH-202621", size=13, color=DIM)
    tb(s, 1.0, 1.25, 11.5, 1.2, "多源异构数据驱动岗位和能力图谱构建\n与动态演化分析研究", size=34, bold=True, spacing=1.15)
    rect(s, 1.0, 3.25, 1.6, 0.03, fill=ACCENT)
    tb(s, 1.0, 3.45, 11.5, 0.9, [
        {"runs": [("智演 ", {"size": 30, "bold": True}), ("JobEvolution", {"size": 30, "bold": True, "color": ACCENT})]},
        {"text": "从 53 个企业官方招聘门户长出可对账的岗位能力图谱：发现新岗位、记录既有岗位要求怎么变、对着一份简历给出方向结论",
         "size": 14, "color": DIM, "after": 0},
    ], spacing=1.2)
    # stat row
    xs = [1.0, 3.6, 6.2, 8.8]
    for (x, n, l) in zip(xs, ["53", "3,580", "4 / 17", "≥90%"],
                          ["企业官方招聘门户", "去重 JD 快照", "领域 / 规范岗位靶子", "简历提取与匹配 F1 口径"]):
        stat(s, x, 5.05, 2.4, n, l, num_size=30)
    tb(s, 1.0, 6.45, 11, 0.8, [
        {"text": "参赛单位：__________________　　团队成员：__________________　　指导教师：__________________", "size": 12, "color": INK},
        {"text": "2026 年 9 月", "size": 12, "color": DIM},
    ])
    s.notes_slide.notes_text_frame.text = (
        "各位评委好。我们的作品叫智演，对应科大讯飞发榜的 XH-202621 赛题。"
        "一句话：它从企业官方招聘门户的 JD 里，自动长出一张能点到证据的岗位能力图谱，能发现新岗位、记录既有岗位要求的变化，"
        "并对着一份简历给出方向结论。接下来 10 分钟按问题、方案、四项功能、创新点、可验证性、提交物的顺序介绍。")


# =============================================================== 2 problem
def s_problem():
    s = chrome("01 问题", "岗位在变，但求职者拿到的是一份静态描述",
               "劳动力市场的岗位名和技能点每季度都在变；现有做法既跟不上，也不可信",
               notes="先说问题。第一，时滞：招聘网站上的岗位描述是抄来抄去的静态文本，Agent 工程师已经从大模型应用工程师里拆出来了，"
                     "但没有一个地方能告诉你它什么时候成为独立岗位、要求什么。第二，噪声和抄袭：同一份 JD 被多个渠道转发，"
                     "统计时被重复计数，通用素质和模型品牌被当成技能。第三，幻觉：直接让大模型写岗位画像，它会编出不存在的要求，"
                     "求职者无法核对。我们把这三个问题作为设计起点。")
    cards = [
        ("时滞", "岗位名称与要求每季度变化，人工维护的岗位字典滞后 6–12 个月；\n新岗位（如 Agent 工程师）何时从既有岗位拆出，没有可追溯的判定依据。"),
        ("噪声与抄袭", "同一 JD 被多渠道转发、模板化抄写，重复计票放大伪信号；\n“沟通能力”“GPT-4”等通用素质与模型品牌混入技能点，污染要求边。"),
        ("幻觉", "直接让大模型生成岗位画像或匹配报告，结论无法回到原文；\n一个 0–100 的匹配分看起来精确，却没人能解释它从哪里来。"),
    ]
    x = 0.6
    for head, body in cards:
        card(s, x, 2.05, 3.9, 2.55, head, body, size=12.5, head_size=16)
        x += 4.12
    tb(s, 0.6, 4.85, 12.1, 0.4, "我们的设计目标", size=13, bold=True, color=ACCENT)
    tb(s, 0.6, 5.2, 12.1, 1.6, [
        {"text": "每一条岗位要求都能点到原文证据和独立来源；拿不出证据的不进图谱。", "bullet": True},
        {"text": "岗位状态、技能入池、要求性质、匹配档位全部由确定性规则从证据计票得出；大模型只负责“读”，不负责“算”。", "bullet": True},
        {"text": "求职者看到的是可核对的方向结论、缺口与换档条件，不是一个大数字。", "bullet": True},
    ], size=13.5, spacing=1.25, space_after=6)


# =============================================================== 3 requirement mapping
def s_mapping():
    s = chrome("02 赛题对应", "赛题四项核心功能与本作品的落点",
               notes="这一页把赛题要求和我们的实现逐条对上。四项核心功能分别落在市场变化页、岗位工作台、图谱本体和五步诊断上；"
                     "创新要求里的多源交叉验证和幻觉防控是贯穿管线的闸门，不是单独的功能页。可验证性要求我们在后面单独讲。")
    rows = [
        ["赛题要求", "本作品实现", "落点"],
        ["① 新岗位发现与定义", "未对齐 17 个规范岗位靶子的 JD → bge 嵌入层次聚类 → LLM 三分类（新岗位 / 别名 / 噪声）→ 状态机 候选→萌芽→成型；岗位定义带证据、需审批", "/discover 市场变化卷宗"],
        ["② 既有岗位能力动态更新", "每周期重算簇内覆盖率（≥30% 入池）与要求判定票（明确票 ≥60% ∧ ≥2 独立源）；旧边写 valid_to 不删除，切片差分标新增 / 升值 / 失效", "/graph 岗位工作台 · /admin 待审"],
        ["③ 全景图谱到技能点粒度", "六类节点（领域 / 岗位 / 技能类目 / 技能点 / 证据 / 演化事件）；按领域、技术栈类目、熟练级切换；必备 / 加分 / 观测中分组", "/graph 图谱与要求清单双视图"],
        ["④ 人岗匹配诊断与差距分析", "PDF/docx 简历解析并校对 → 岗位推荐序 → 双岗确定性对照（档位 / 缺口集 / 换档条件）→ 模型解释每条带证据 ID → 双轨行动清单", "/diagnose 五步诊断"],
        ["创新：多源清洗与交叉验证", "fingerprint 幂等 · simhash 近重 · 独立源 = 规范化公司 · 双时间戳 · 通用素质与模型品牌出边", "collectors / pipeline"],
        ["创新：幻觉防控", "原文回指 → 三层归一 → 入池门 → 判定票门 → 置信层 → 审核闸 → 发布校验，共十道闸", "pipeline/gate · diagnostic_release"],
        ["可验证性", "100 条 JD / 100 份简历 / 100 对匹配金标；set-based F1；pytest-cov 门槛 60%；Docker Compose 一键部署", "data/eval-official-only · tests/"],
    ]
    table(s, 0.6, 2.0, 12.1, rows, [2.6, 6.9, 2.6], size=11, row_h=0.56, first_bold=True)


# =============================================================== 4 architecture
def s_arch():
    s = chrome("03 总体设计", "六层架构：数据从官方门户向上长成图谱，再服务于诊断",
               notes="总体架构分六层。最底层是 53 个企业官方招聘门户；采集层做指纹幂等和近重去重；"
                     "图谱构建层是核心，包含切段、LLM JSON 抽取、三层归一、入池与判定票、置信层与审核闸、发布校验；"
                     "存储用 Neo4j 存图、Redis 存会话和事件流；服务层是 FastAPI；应用层是 Next.js 的五个页面。"
                     "大模型的唯一出口是 app/llm/client.py，方便更换供应商和 mock 测试。")
    pic(s, f"{DIAG}/fig_arch.png", 0.6, 1.95, 7.6, 5.0, mode="contain", border=False)
    tb(s, 8.5, 2.0, 4.3, 4.9, [
        {"text": "技术选型", "size": 14, "bold": True, "color": ACCENT, "after": 8},
        {"text": "前端  Next.js 15 · React 19 · TypeScript；React Flow + d3-force（Web Worker 布局）", "bullet": True},
        {"text": "服务  Python 3.12 · FastAPI · Pydantic；SSE 推送审核与管线事件", "bullet": True},
        {"text": "存储  Neo4j 5 Community（图谱、证据、演化事件、发布版本）· Redis 7 AOF（简历会话 TTL 1h、指纹、事件流）", "bullet": True},
        {"text": "模型  OpenAI 兼容接口（DeepSeek / B.AI / Tuzi），JSON 模式、非思考；嵌入 BAAI/bge-m3", "bullet": True},
        {"text": "部署  Docker Compose 五服务 + 测试库 profile；GitHub Actions CI 跑 mock 评测与覆盖率门槛", "bullet": True},
        {"text": "原则  大模型只“读”不“算”；每条要求边可回到原文；旧事实写失效时间不删除", "bullet": True, "bold": True},
    ], size=12, spacing=1.22, space_after=6)


# =============================================================== 5 data sources
def s_data():
    s = chrome("04 多源异构数据", "53 个企业官方门户 · 3,580 条去重 JD · 四类问题各有一把确定性的刀",
               notes="数据来源是 53 个企业官方招聘门户，覆盖飞书、腾讯、字节、北森四种 ATS 结构，不抓第三方招聘网站，避免二次转发的抄袭。"
                     "清洗做四件事：正文 sha256 指纹保证幂等；64 位 simhash 海明距离不超过 3 判近重，只留最早一条；"
                     "独立源按规范化公司名去重，一家公司多渠道只算一票；观察时间和有效期分开记，这样切片差分才能重演。"
                     "抽取阶段只允许职责段和要求段进入技能点，福利和公司介绍不计。")
    xs = [0.6, 3.6, 6.6, 9.6]
    for x, n, l in zip(xs, ["53", "4 类", "3,580", "≤3"],
                        ["启用的官方招聘门户（飞书 / 腾讯 / 字节 / 北森）", "ATS 门户适配器，统一 JD 快照结构", "去重后的 JD 快照，双时间戳落盘", "simhash 64 位海明距离阈值，近重只留最早"]):
        stat(s, x, 1.95, 2.8, n, l, num_size=30)
    rows = [
        ["问题", "确定性对策", "实现位置"],
        ["时滞", "每日增量采集，连续两页指纹未变即提前停止；observed_at 与 valid_from / valid_to 分开；本周期切片差分直接标新增 / 升值 / 失效", "collectors · Evidence.observed_at"],
        ["噪声", "只取职责段与要求段；通用素质（沟通、团队协作）不进技能点；模型品牌默认留在证据中，从同句动作抽 API 集成 / 部署 / 微调 / 评测", "pipeline/sections · extract 提示词"],
        ["抄袭 / 重复", "正文 sha256 fingerprint 幂等；simhash 近重不计票；独立源 = 规范化公司名，渠道不计票", "collectors/dedupe · gate.summarize_requirement_votes"],
        ["通胀", "斜杠并列拆开、原文“任选 / 或”建要求组读 min_required；发布上限 12 必备 / 24 正式；新增超 50% 记岗位要求异常并暂停诊断", "pipeline/diagnostic_release"],
    ]
    table(s, 0.6, 3.45, 12.1, rows, [1.4, 7.6, 3.1], size=11, row_h=0.6, first_bold=True)


# =============================================================== 6 new job discovery
def s_discover():
    s = chrome("05 功能一 · 新岗位发现与定义", "候选 → 萌芽 → 成型：岗位状态由独立源计票自动流转，不手写",
               notes="第一项功能。对不上 17 个规范岗位靶子的 JD，用标题加技能点做 bge 嵌入，层次聚类，最小簇 3；"
                     "再让大模型对每个簇做三分类：新岗位、别名、噪声。别名并入既有岗位，比如“大模型应用开发工程师”判为大模型应用工程师的别名。"
                     "状态机：候选不入谱；3 个独立源、90 天窗、判为新岗位才成为萌芽；10 个独立源或持续 6 个月且岗位定义获批才成型。"
                     "右边是市场变化页，用求职者的问题组织卷宗：为什么系统认为岗位正在形成、哪些公司在招、与最接近的既有岗位有什么区别。")
    pic(s, f"{DIAG}/fig_states.png", 0.6, 1.95, 6.9, 2.85, mode="contain", border=False)
    tb(s, 0.6, 4.9, 6.9, 2.0, [
        {"text": "赛题示例：Agent 工程师", "size": 13, "bold": True, "color": ACCENT, "after": 6},
        {"text": "从大模型应用工程师簇中分离；独立源 ≥3 后进入萌芽，卷宗记录首次出现时间、招聘公司、本周期新增要求与相近岗位差异。", "bullet": True},
        {"text": "「大模型应用开发工程师」由 LLM 簇判别为别名，写 ALIAS_OF 并入既有岗，不占候选列。", "bullet": True},
        {"text": "岗位定义（一句定义 + 3–5 条典型职责）由模型起草、每条带证据，管理员批准后才公开；未获批显示“岗位定义审核中”。", "bullet": True},
    ], size=12, spacing=1.2, space_after=5)
    pic(s, f"{SHOT}/discover.png", 7.8, 1.95, 4.9, 3.06, mode="cover")
    tb(s, 7.8, 5.1, 4.9, 0.3, "市场变化页：候选 / 萌芽 / 成型筛选，公司按规范名去重", size=10, color=DIM)
    tb(s, 7.8, 5.5, 4.9, 1.4, [
        {"text": "阈值（pipeline/constants.py）", "size": 12, "bold": True, "after": 4},
        {"text": "DISCOVER_MIN_CLUSTER 3 · EMERGING_SOURCES 3 · EMERGING_WINDOW_DAYS 90 · FORMED_SOURCES 10 · FORMED_MONTHS 6 · JOB_ALIGN_THRESHOLD 0.84", "size": 11, "color": DIM, "font": "Menlo"},
    ], size=12)


# =============================================================== 7 dynamic update
def s_update():
    s = chrome("06 功能二 · 既有岗位能力动态更新", "每份去重 JD 只投一张票；入池看覆盖率，要求性质看判定票",
               notes="第二项功能。技能点要进入岗位要求，先过入池门：来自职责或要求段，且簇内覆盖率不低于 30%，否则只是“观测中”。"
                     "然后每份去重 JD 投一张要求判定票：原文写“必须、熟练掌握”记明确必备，写“优先、加分”记明确加分，普通提及记未标。"
                     "明确票占已分类票 60% 以上，且来自至少两个独立源，才能确定要求性质；否则只进待审。"
                     "旧的要求边不删，写 valid_to，所以切片差分可以重演。右边是管理后台，按岗位与待发布版本分组，"
                     "管理员能看到三类票数和原文，批、改、驳都留痕。")
    tb(s, 0.6, 1.95, 5.6, 4.9, [
        {"text": "三道门", "size": 14, "bold": True, "color": ACCENT, "after": 6},
        {"text": "入池门  职责 / 要求段 ∧ 簇内覆盖率 ≥30%；否则标“观测中”，不是缺口", "bullet": True},
        {"text": "判定票门  明确必备或明确加分票 ≥60% 已分类票 ∧ ≥2 独立源；否则只生成提案", "bullet": True},
        {"text": "置信层  高 / 中 / 低；低层永不自动入谱，自动审核开关默认关，开启后仍需独立审核模型复核", "bullet": True},
        {"text": "切片差分", "size": 14, "bold": True, "color": ACCENT, "after": 6},
        {"text": "本周期覆盖率上升并刚过入池线 → 升值；失效写 valid_to → 可后补；新增 / 升值 / 失效三组文字清单，不让用户从画布猜", "bullet": True},
        {"text": "发布校验", "size": 14, "bold": True, "color": ACCENT, "after": 6},
        {"text": "定义非空、证据存在且未撤回、无重复有效要求、要求等价数 ≤12 必备 / ≤24 正式；新增超 max(3, 上期×50%) 记异常，只暂停该岗诊断", "bullet": True},
        {"text": "示例  大模型应用工程师 2026-09-02 校准发布：12 条必备 / 23 条正式；通用素质与模型品牌以 valid_to 失效、可回滚", "size": 11.5, "color": DIM},
    ], size=12, spacing=1.2, space_after=5)
    pic(s, f"{SHOT}/admin-queue.png", 6.5, 1.95, 6.2, 3.875, mode="cover")
    tb(s, 6.5, 5.9, 6.2, 0.9, [
        {"text": "管理后台：按岗位与待发布版本分组，展示上期 / 本期要求等价数、三类要求判定票与最短原文；一键全部批准仍执行确定性校验，审计记录保存会话、时间、提案 ID 与放行理由。", "size": 10.5, "color": DIM},
    ])


# =============================================================== 8 graph
def s_graph():
    s = chrome("07 功能三 · 全景图谱可视化", "六类节点到技能点粒度，按领域 / 技术栈 / 熟练级切换",
               "每个技能点可打开证据抽屉；级别写在要求边上，岗位不按初高级分裂",
               notes="第三项功能，全景图谱。本体有六类节点：领域、岗位、技能类目、技能点、证据、演化事件。"
                     "要求边上带必备或加分、熟练级、有效期、置信层、证据 ID、要求组信息，所以级别写在边上，岗位不按初高级分裂。"
                     "工作台画当前岗位的切片，用领域、技术栈类目、熟练级筛选；桌面默认图谱视图、移动端默认要求清单，两者随时切换。"
                     "点任何技能点都能打开证据抽屉看最短原文和公司来源。图布局在 Web Worker 里跑 d3-force，主线程不卡。")
    pic(s, f"{SHOT}/graph.png", 0.6, 1.95, 7.6, 4.75, mode="cover")
    tb(s, 0.6, 6.75, 7.6, 0.3, "岗位工作台 /graph：岗位说明区 · 领域 / 类目 / 级别筛选 · 图谱 / 要求清单切换 · 切片差分与判定摘要", size=10, color=DIM)
    pic(s, f"{DIAG}/fig_ontology.png", 8.45, 1.95, 4.3, 2.2, mode="contain", border=False)
    tb(s, 8.45, 4.25, 4.3, 2.7, [
        {"text": "本体六类：Domain · Job · SkillCategory · Skill · Evidence · EvolutionEvent", "bullet": True},
        {"text": "REQUIRES 边属性：kind 必备/加分 · proficiency 熟练级 · valid_from/valid_to · confidence · layer · sources[] · group_id/min_required", "bullet": True},
        {"text": "切换维度：领域 → 岗位 → 技能类目（语言 / 框架 / 平台 / 工程 / 领域知识）→ 熟练级（aware / able / expert）", "bullet": True},
        {"text": "求职者文案用人话：独立源 = “去重招聘公司”，切片 = “截至某日的岗位要求”", "bullet": True},
    ], size=11.5, spacing=1.2, space_after=5)


# =============================================================== 9 diagnose flow
def s_diagnose():
    s = chrome("08 功能四 · 人岗匹配诊断与差距分析", "确定性匹配算档位与缺口，大模型只解释、每条判断带证据 ID",
               notes="第四项功能，五步诊断。上传 PDF 或 docx 简历，只取文本层，原文件即删，会话存 Redis 一小时，免登录。"
                     "模型解析后用户按原文校对技能点、熟练级和证据级。系统按岗位推荐序给出三个可诊断岗位，用户选两个对照。"
                     "匹配分是确定性公式：必备覆盖加 0.3 倍加分覆盖，半档记 0.5，四个档位阈值 0.85、0.60、0.35。"
                     "大模型只做方向结论的解释，不合成新分数，每条判断带证据 ID，缺失只能写“简历中未找到”。")
    pic(s, f"{DIAG}/fig_diagnose.png", 0.6, 1.9, 7.3, 2.95, mode="contain", border=False)
    card(s, 8.1, 1.95, 4.6, 2.75, "确定性计分",
         [{"text": "匹配分 = 100 × (必备覆盖 + 0.3 × 加分覆盖) / (必备满分 + 0.3 × 加分满分)", "size": 11.5},
          {"text": "熟练级低一级记半档 0.5；加分缺失不伤必备；经验、学历只作旁注不进档", "size": 11.5},
          {"text": "档位：≥0.85 高度匹配 · ≥0.60 基本匹配 · ≥0.35 有明显差距 · 其余不匹配", "size": 11.5}], size=11.5)
    card(s, 0.6, 4.95, 5.95, 1.95, "缺口集与换档条件",
         [{"text": "缺口 = 已发布必备要求中简历无证据或熟练级不足者；要求组按 min_required 对账，不重复制造缺口", "size": 11.5},
          {"text": "换档条件 shift_set：单独一项能升档者排前，成对才能换档者次之；学习路径按此排序", "size": 11.5}], size=11.5)
    card(s, 6.75, 4.95, 5.95, 1.95, "简历证据级",
         [{"text": "提及（技能栏 / 自我评价）· 使用（经历写明任务中的使用）· 结果（本人责任 + 数字 / 交付物）", "size": 11.5},
          {"text": "最高只有提及记“简历证据不足”；证据级不进匹配分，只进方向结论与改写建议", "size": 11.5}], size=11.5)


# =============================================================== 10 report
def s_report():
    s = chrome("09 诊断报告", "首屏一句方向结论；结论 / 行动 / 依据三视图；四个产品记忆点",
               notes="这是报告页。首屏固定一句方向结论，下面只放当前档位、最大优势和最大阻碍。"
                     "当两个岗位的档位、必备覆盖、专属技能数、可迁移能力数和最小换档等价数全部相同，系统写“当前简历证据不足以明确区分”，不强选胜者。"
                     "四个记忆点：换档模拟器只接受缺口和熟练级不足项，用现有规则即时重算，结果标“假设结果”；"
                     "简历证据地图左边经历右边要求，连线标提及、使用或结果；邻近岗位迁移地图给出最小换档技能数；"
                     "市场信号雷达列观测中技能的覆盖率和来源，明确说“市场开始提，但还没进要求，不算你的缺口”。")
    pic(s, f"{SHOT}/diagnose-report-detail3.png", 0.6, 1.95, 7.2, 4.38, mode="cover")
    tb(s, 0.6, 6.4, 7.2, 0.5, "单岗结论视图：方向结论 → 简历定位判断 → 关键优势 / 主要风险（每条可打开简历证据片段）→ 档位与换档条件", size=10, color=DIM)
    items = [
        ("换档模拟器", "只接受缺口、熟练级不足、要求组候选技能；即时重算两岗与邻近岗；固定标“假设结果，尚未被简历证明”"),
        ("简历证据地图", "左列经历 / 项目证据片段，右列正式要求，连线标提及 / 使用 / 结果；未连接的要求写“简历未找到证据”"),
        ("邻近岗位迁移地图", "两个对照岗 + 推荐序中一个邻近岗；逐岗给当前档位、最小换档技能数、共享能力、独有要求"),
        ("市场信号雷达", "观测中技能按覆盖率、独立源、周期变化列出，并说明未进要求的原因；不算缺口"),
    ]
    y = 1.95
    for head, body in items:
        rect(s, 8.1, y, 4.6, 1.02, fill=CREAM)
        rect(s, 8.1, y, 0.06, 1.02, fill=ACCENT)
        tb(s, 8.3, y + 0.08, 4.3, 0.3, head, size=12.5, bold=True)
        tb(s, 8.3, y + 0.4, 4.3, 0.62, body, size=10.5, color=INK, spacing=1.15)
        y += 1.12
    tb(s, 8.1, 6.45, 4.6, 0.5, "双轨行动清单：简历证明轨（五条优先改写，原文 → 问题 → 建议版本 → 待补事实）+ 能力提升轨（≤3 项，按换档条件）", size=10, color=DIM)


# =============================================================== 11 hallucination gates
def s_gates():
    s = chrome("10 创新一 · 幻觉防控", "十道闸门：未过只降级，不进正式图谱",
               "任何一道未过只降级为观测中 / 审核提案 / 暂停诊断；大模型输出从不直接成为事实",
               notes="第一个创新点是幻觉防控，这也是赛题明确要求的。我们不靠提示词祈祷，而是把大模型输出放进十道确定性闸门。"
                     "抽取时每一项必须带原文摘录，摘录在原文里命不中直接丢弃；技能名走三层归一，跨语言或缩写关系只能生成合并提案由人核对；"
                     "入池门、判定票门看的是覆盖率和独立源票数，模型改不了；置信层低的永不自动入谱；审核闸默认人审，"
                     "自动审核开启也要独立模型复核；发布校验守数量上限和异常增长；最后发布版本是不可变事实集，可以回滚。"
                     "诊断侧同样：匹配分由规则算，模型解释每条带证据 ID，缺失只能写“未找到”。")
    pic(s, f"{DIAG}/fig_gates.png", 0.6, 1.9, 12.1, 3.2, mode="contain", border=False)
    cols = [
        ("构建侧", ["原文回指：摘录必须在 JD 原文命中", "三层归一：表面归一 → 获批同义词 → bge 余弦 ≥0.70；相关技术禁止合并", "入池门 30% · 判定票门 60% ∧ 2 独立源", "置信层：低层永不自动入谱", "审核闸：人工批 / 改 / 驳；自动审核需独立模型复核"]),
        ("发布侧", ["定义非空、证据存在且未撤回、无重复有效要求", "等价数上限 12 必备 / 24 正式", "新增超 max(3, 上期×50%) → 岗位要求异常，暂停该岗诊断", "发布版本不可变，可回滚；撤回 ≠ 演化"]),
        ("诊断侧", ["匹配分、档位、缺口集由 bands.py 确定性计算", "模型解释不合成分数，每条判断带证据 ID", "缺失项只能写“简历中未找到”，并标明检查了哪个部分", "改写建议不得生成原文没有的数字与结果，只给待补占位"]),
    ]
    x = 0.6
    for head, lines in cols:
        tb(s, x, 5.15, 3.95, 1.8, [{"text": head, "size": 12, "bold": True, "color": ACCENT, "after": 2}] +
           [{"text": l, "bullet": True, "size": 10, "after": 0} for l in lines], spacing=1.08)
        x += 4.08


# =============================================================== 12 innovation 2
def s_innov2():
    s = chrome("11 创新二 · 多源交叉验证与“读算分离”", "大模型只负责读，图谱与规则负责算；结论必须能回到证据与来源",
               notes="第二个创新点讲大模型、图谱、检索三者的协同方式，我们叫读算分离。"
                     "大模型做四件“读”的事：JSON 抽取、簇判别、岗位定义起草、诊断解释。规则和图谱做所有“算”的事：覆盖率、判定票、状态机、匹配分、换档条件。"
                     "交叉验证体现在独立源：一个结论要成立，必须有至少两家不同公司的 JD 支持，同一公司多渠道不重复计票。"
                     "检索增强不是把 JD 喂进上下文让模型总结，而是把图谱里已发布的要求边、证据和简历证据片段作为结构化输入，"
                     "模型只能在这些 ID 上引用。这让每条输出都能被评委和用户点开核对。")
    card(s, 0.6, 1.95, 3.9, 2.4, "大模型只“读”", [
        {"text": "JSON 模式抽取技能点，每项带原文摘录与段落来源", "bullet": True},
        {"text": "对聚类簇做三分类：新岗位 / 别名 / 噪声", "bullet": True},
        {"text": "起草岗位定义与金标草稿（ADR-0011：LLM 起草、人裁决）", "bullet": True},
        {"text": "诊断解释：方向结论、优势风险、改写建议", "bullet": True},
    ], size=11.5)
    card(s, 4.7, 1.95, 3.9, 2.4, "图谱与规则“算”", [
        {"text": "簇内覆盖率、要求判定票、独立源计数", "bullet": True},
        {"text": "岗位状态机与 ALIAS_OF 并入", "bullet": True},
        {"text": "匹配分、档位、缺口集、换档条件 shift_set", "bullet": True},
        {"text": "发布校验与岗位要求异常", "bullet": True},
    ], size=11.5)
    card(s, 8.8, 1.95, 3.9, 2.4, "结构化检索增强", [
        {"text": "解释阶段输入的是已发布要求边、证据 ID 与简历证据片段 ID，不是整篇 JD", "bullet": True},
        {"text": "模型只能引用给定 ID；未引用的判断被丢弃", "bullet": True},
        {"text": "学习资源与岗位职责摘录一起给出“为何先补它”", "bullet": True},
    ], size=11.5)
    tb(s, 0.6, 4.6, 12.1, 0.4, "多源交叉验证在三个层面生效", size=13, bold=True, color=ACCENT)
    rows = [
        ["层面", "规则", "解决的问题"],
        ["证据层", "fingerprint 幂等 · simhash 近重只留最早 · 只取官方门户不取转发站", "抄袭、重复计数"],
        ["岗位层", "独立源 = 规范化公司名；萌芽 ≥3 源、成型 ≥10 源或 ≥6 月；别名并入不占候选", "伪新岗位、名称漂移"],
        ["要求层", "明确票 ≥60% ∧ ≥2 独立源；覆盖率随周期重算；合并提案人工核对", "单家公司偏好被当成市场要求；同义词漂移"],
    ]
    table(s, 0.6, 5.05, 12.1, rows, [1.6, 7.2, 3.3], size=11, row_h=0.44, first_bold=True)


# =============================================================== 13 verification
def s_verify():
    s = chrome("12 可验证性", "三套金标、set-based F1、覆盖率门槛与金标修订规程",
               notes="可验证性。我们建了三套各 100 条的金标：JD 解析、简历提取、匹配缺口集，全部用图谱 Skill.id 标注，"
                     "对齐阈值冻结在 freeze.json 里，评测只读这个文件。金标由大模型起草、人逐条回看原文裁决，修订走两段规程，禁止照着系统预测改金标。"
                     "最近一次未 mock 结果：简历提取和匹配都是 1.000，JD 解析 0.814 还没到 0.90，差距样本和复跑要求写在 summary.md 里，我们没有改金标凑数。"
                     "单测 166 通过，覆盖率 70.69%，高于 60% 门槛，核心管线模块覆盖率都在 80% 以上。")
    rows = [
        ["指标", "赛题要求", "数据集", "最近一次未 mock 结果", "状态"],
        ["JD 解析 F1", "≥0.90", "jd.jsonl 100 条（去重、四领域、17 岗各 ≥3）", "0.814", "未达标，差距样本已归档"],
        ["简历提取 F1", "≥0.90", "resume.jsonl 100 份（单栏 / 双栏、PDF / docx）", "1.000", "达标"],
        ["匹配 F1（缺口集）", "≥0.90", "match_pairs.jsonl 100 对（金标简历 × 目标岗）", "1.000", "达标"],
        ["单元测试", "覆盖率 ≥60%", "pytest --cov（分母排除 collectors、llm/client 网络分支）", "166 passed · 70.69%", "达标"],
        ["JD 测试集规模", "≥100 条去重", "simhash 近重只留最早；官方门户快照", "100 条", "达标"],
    ]
    table(s, 0.6, 1.95, 12.1, rows, [1.9, 1.5, 4.4, 2.3, 2.0], size=11, row_h=0.46, first_bold=True)
    tb(s, 0.6, 4.9, 6.0, 2.0, [
        {"text": "口径与规程", "size": 13, "bold": True, "color": ACCENT, "after": 5},
        {"text": "预测集与金标集都先过 align_skill（读 freeze.json），对技能点集合算 P / R / F1；空对空记 1，一边空记 0", "bullet": True},
        {"text": "匹配评测不喂解析输出：系统用金标简历技能 × 金标要求边再算缺口集，与人工缺口集比集合", "bullet": True},
        {"text": "金标修订两段：先盲改只看原文与词表，再用系统预测找分歧逐条回原文裁决，理由写进 notes（ADR-0011）", "bullet": True},
        {"text": "CI 跑 mock 三项 + 覆盖率门槛；summary.md 数字只来自本地未 mock 跑", "bullet": True},
    ], size=11, spacing=1.18, space_after=3)
    tb(s, 6.9, 4.9, 5.8, 2.0, [
        {"text": "模块覆盖率（本次 pytest-cov）", "size": 13, "bold": True, "color": ACCENT, "after": 5},
        {"text": "sections 100% · score 98% · bands 97% · extract 94% · diagnostic_release 93% · curate_public 89% · status 88% · resume 85% · gate 81% · graph 77% · align 71%", "size": 11},
        {"text": "JD 0.814 的差距来源：模型输出截断与词表边界样本；改进计划：分段重试上限、候选召回只接受精确词表命中、复跑要等 100 条完整样本", "size": 11, "color": DIM},
    ], size=11, spacing=1.18)


# =============================================================== 14 implementation & deploy
def s_deploy():
    s = chrome("13 系统实现与部署", "Docker Compose 一键启动五个服务；测试库独立卷，永不写产品图",
               notes="系统实现。仓库是 monorepo：apps/api 是 FastAPI，apps/web 是 Next.js，data 目录放 JD 快照、金标和图谱快照。"
                     "docker compose up 一键起 web、api、pipeline、neo4j、redis 五个服务，pipeline 容器每日采集加抽取。"
                     "测试用独立的 neo4j-test profile，pytest 固定连测试库，连不上直接报错，保证永远不会污染产品图。"
                     "空卷导入快照几分钟就能复现完整图谱，评委可以在自己机器上跑起来。")
    pic(s, f"{DIAG}/fig_deploy.png", 0.6, 1.95, 7.0, 2.75, mode="contain", border=False)
    tb(s, 0.6, 4.85, 7.0, 2.1, [
        {"text": "一键启动", "size": 13, "bold": True, "color": ACCENT, "after": 4},
        {"text": "cp .env.example .env   # 填 LLM_PROVIDER / API KEY / ADMIN_PASSWORD\ndocker compose up -d --build\ndocker compose --profile test up -d neo4j-test && pytest --cov -q", "size": 11, "font": "Menlo", "color": INK},
        {"text": "健康检查：GET /v1/meta 返回图谱发布版本、模型供应商、简历保留策略；Neo4j 7474 与 redis-cli ping", "size": 11, "color": DIM},
    ], size=11, spacing=1.2)
    rows = [
        ["路由", "职责"],
        ["GET /jobs · /jobs/{id} · /graph", "岗位列表、岗位切片与差分、图谱节点边"],
        ["GET /discover · /feed", "候选 / 萌芽 / 成型卷宗、总览故事流"],
        ["POST /sessions · PATCH /sessions/{id}", "简历上传解析、校对修正（Redis TTL 1h）"],
        ["POST /diagnose · /diagnose/simulate", "双岗对照报告、换档模拟"],
        ["/admin/* · SSE /events", "口令门、待审队列、批 / 改 / 驳、自动审核开关、审计"],
        ["GET /v1/meta", "图谱发布版本、供应商、运行状态"],
    ]
    table(s, 7.9, 1.95, 4.8, rows, [2.3, 2.5], size=10.5, row_h=0.42, first_bold=True)
    tb(s, 7.9, 5.05, 4.8, 1.9, [
        {"text": "代码组织", "size": 13, "bold": True, "color": ACCENT, "after": 4},
        {"text": "apps/api/app/{collectors, pipeline, matching, eval, llm, graph}", "size": 11, "font": "Menlo"},
        {"text": "apps/web/app/{graph, discover, diagnose, admin} + components", "size": 11, "font": "Menlo"},
        {"text": "data/{official-only, eval-official-only, snapshot}", "size": 11, "font": "Menlo"},
        {"text": "docs/{product, tech, verification, adr/}  ·  CONTEXT.md 术语权威", "size": 11, "font": "Menlo"},
    ], size=11, spacing=1.2, space_after=2)


# =============================================================== 15 deliverables
def s_deliver():
    s = chrome("14 提交物清单", "赛题要求的七项提交物与对应位置",
               notes="提交物一共七项。方案文档和 PPT 就是现在这两份；演示视频不超过 10 分钟，按下一页分镜录制；"
                     "源码含 Docker Compose 部署脚本；单元测试和覆盖率报告可以一条命令复现；"
                     "两岗测试数据在 deliver 目录：Agent 工程师作为新岗位，大模型应用工程师作为既有岗位，各含岗位定义、去重证据和输入输出说明。")
    rows = [
        ["#", "提交物", "内容 / 位置", "状态"],
        ["1", "作品设计实现方案", "XH-202621_作品设计实现方案_智演JobEvolution.docx（15 章 + 2 附录，14 图 21 表）", "已完成"],
        ["2", "作品 PPT", "本文件；含每页讲稿备注，对应 10 分钟视频分镜", "已完成"],
        ["3", "演示视频 ≤10 分钟", "按下一页分镜录制：问题 → 图谱 → 发现 → 诊断 → 后台 → 验证", "待录制"],
        ["4", "源码", "GitHub 仓库 JobEvolution（apps/api · apps/web · data · docs · tests）", "已完成"],
        ["5", "Docker 部署", "docker-compose.yml 五服务 + neo4j-test profile；.env.example；快照空卷导入", "已完成"],
        ["6", "单元测试", "tests/ · pytest --cov 166 passed · 70.69%（门槛 60%）· GitHub Actions CI", "已完成"],
        ["7", "两岗测试数据", "data/eval-official-only/deliver/agent（新岗位）· deliver/llm-app（既有岗位）：job.json · sources.jsonl · io.md · diagnose.example.json", "已完成"],
    ]
    table(s, 0.6, 1.95, 12.1, rows, [0.5, 2.4, 7.6, 1.6], size=11, row_h=0.5, first_bold=True)
    tb(s, 0.6, 6.15, 12.1, 0.8, [
        {"text": "评测数据：data/eval-official-only/{jd, resume, match_pairs}.jsonl 各 100 条 · freeze.json（对齐阈值、模型、日期）· out/summary.md（三项 F1、覆盖率、freeze 哈希）", "size": 11, "color": DIM},
    ])


# =============================================================== 16 video storyboard
def s_video():
    s = chrome("15 演示视频分镜", "10 分钟：问题 1 分 · 系统 6 分 · 验证 2 分 · 收尾 1 分",
               notes="这一页是我们的录制分镜，也是答辩节奏。前一分钟讲问题；中间六分钟走一遍系统：总览、图谱工作台看切片差分、"
                     "市场变化页看 Agent 工程师卷宗、上传简历走完五步诊断看报告四个记忆点、管理后台看判定票和审核；"
                     "两分钟讲评测口径与结果、一条命令复现；最后一分钟总结创新点与价值。")
    rows = [
        ["时间", "镜头", "画面", "要说的一句话"],
        ["0:00–1:00", "问题", "首页第一屏 + 三个痛点", "岗位在变，静态描述跟不上，大模型直接生成又不可信"],
        ["1:00–2:00", "总览与图谱", "/ 故事流 → /graph 岗位工作台", "每条要求边能点到证据；切片差分标出本周期新增 / 升值 / 失效"],
        ["2:00–3:00", "新岗位发现", "/discover Agent 工程师卷宗", "3 个独立源、90 天、判为新岗位才成萌芽；别名并入不占位"],
        ["3:00–5:30", "五步诊断", "上传 → 校对 → 推荐序选两岗 → 报告", "确定性档位与缺口；换档模拟器、证据地图、迁移地图、市场雷达"],
        ["5:30–6:30", "既有岗位更新", "/graph 要求清单判定摘要 → 判断依据", "明确票 ≥60% ∧ 2 独立源；观测中不是缺口"],
        ["6:30–7:00", "管理后台", "/admin 待审、票数、原文、一键批准", "待审是闸不是第二套产品；审计留痕"],
        ["7:00–9:00", "可验证性", "终端跑 pytest --cov 与 app.eval；summary.md", "三套金标、set-based F1、覆盖率 70.69%；JD 差距如实归档"],
        ["9:00–10:00", "收尾", "架构图 + 十道闸门图", "读算分离、多源交叉验证、每条结论可回到证据"],
    ]
    table(s, 0.6, 1.95, 12.1, rows, [1.4, 1.8, 3.9, 5.0], size=11, row_h=0.5, first_bold=True)


# =============================================================== 17 closing
def s_close():
    PAGE[0] += 1
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, fill=INK)
    rect(s, 0, 0, 0.35, H, fill=ACCENT)
    tb(s, 1.0, 1.0, 11, 0.4, "总结", size=13, color="c9c9c9")
    tb(s, 1.0, 1.5, 11.5, 2.2, [
        {"runs": [("智演 ", {"size": 34, "bold": True, "color": PAPER}), ("JobEvolution", {"size": 34, "bold": True, "color": ACCENT})]},
        {"text": "一张能点到证据的岗位能力图谱，和一份能被核对的职业方向结论", "size": 18, "color": "e6e6e6"},
    ], spacing=1.2)
    items = [
        ("多源可信", "53 个官方门户 · 指纹幂等 · 近重去重 · 独立源计票 · 双时间戳"),
        ("自动演化", "候选 → 萌芽 → 成型状态机；覆盖率与判定票每周期重算；旧边写 valid_to 可回放"),
        ("读算分离", "大模型只抽取、判别、解释；分数、档位、状态由确定性规则算；十道闸门防幻觉"),
        ("可验证", "三套 100 条金标 · set-based F1 · 覆盖率 70.69% · Docker 一键复现"),
    ]
    y = 3.85
    for head, body in items:
        rect(s, 1.0, y + 0.06, 0.1, 0.1, fill=ACCENT)
        tb(s, 1.3, y - 0.05, 2.2, 0.4, head, size=15, bold=True, color=PAPER)
        tb(s, 3.5, y - 0.05, 9, 0.4, body, size=13, color="d6d6d6")
        y += 0.62
    tb(s, 1.0, 6.5, 11, 0.5, "请评委指正 · 源码与部署脚本见提交物清单", size=12, color="9a9a9a")
    s.notes_slide.notes_text_frame.text = (
        "总结四点：多源可信、自动演化、读算分离、可验证。智演不是又一个岗位百科，而是一张每条边都能点到证据的图谱，"
        "和一份求职者能逐条核对的方向结论。感谢评委，请指正。")


if __name__ == "__main__":
    s_cover(); s_problem(); s_mapping(); s_arch(); s_data(); s_discover(); s_update(); s_graph()
    s_diagnose(); s_report(); s_gates(); s_innov2(); s_verify(); s_deploy(); s_deliver(); s_video(); s_close()
    prs.save(OUT)
    print("saved", OUT, PAGE[0], "slides")
